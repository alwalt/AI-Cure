import os
import time
import tempfile
from io import BytesIO
from pathlib import Path
from typing import List

import pandas as pd
from fastapi import File, Form, HTTPException, Request, UploadFile
from fastapi.responses import JSONResponse
from langchain_community.document_loaders import PyMuPDFLoader
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_community.chat_models import ChatOllama
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import ConversationalRetrievalChain
from langchain.prompts import PromptTemplate
from Bio import SeqIO

# New processing imports
from pptx import Presentation
import pytesseract
from PIL import Image, ImageDraw, ImageFont
from pathlib import Path
from docx import Document as DocxDocument

# Import from their original locations
from utils import clean_dataframe
from config.shared import SESSIONS, initialize_session, chroma_settings, hnsw_metadata

def docx_read(docx_files):
    full_text = ""
    for docx_file in docx_files:
        doc = DocxDocument(docx_file)
        for para in doc.paragraphs:
            full_text += para.text + "\n"
    return full_text.strip()

def image_from_text(slide_content, image_size=(1024, 768), font_size=20):
    image = Image.new("RGB", image_size, "white")
    draw = ImageDraw.Draw(image)
    try:
        font = ImageFont.truetype("arial.ttf", font_size)
    except IOError:
        font = ImageFont.load_default()
    current_y = 10 
    for line in slide_content.split("\n"):
        draw.text((10, current_y), line, fill="black", font=font)
        current_y += font_size + 5
    return image

def extract_text_from_slide(slide):
    text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text + " "
            text += "\n"
        elif shape.has_table:
            for row in shape.table.rows:
                text += "\t".join(cell.text for cell in row.cells) + "\n"
    return text

def pptx_to_images(pptx_path, output_folder):
    prs = Presentation(pptx_path)
    output_folder = Path(output_folder)
    output_folder.mkdir(parents=True, exist_ok=True)
    images = []
    for i, slide in enumerate(prs.slides):
        slide_text = extract_text_from_slide(slide)
        image = image_from_text(slide_text)
        image_path = output_folder / f"slide_{i + 1}.png"
        image.save(image_path)
        images.append(image_path)
    return images

def ocr_images(image_paths):
    all_text = ""
    for i, image_path in enumerate(image_paths):
        try:
            img = Image.open(image_path).convert("RGB")
            ocr_text = pytesseract.image_to_string(img)
            if ocr_text.strip():
                all_text += f"--- Slide {i + 1} ---\n"
                all_text += ocr_text.strip() + "\n\n"
        except Exception as e:
            all_text += f"[OCR error on slide {i + 1}: {e}]\n"
    return all_text.strip()

def pptx_read(pptx_file_path):
    with tempfile.TemporaryDirectory() as temp_dir:
        slide_images = pptx_to_images(pptx_file_path, temp_dir)
        return ocr_images(slide_images)

def fastq_read(fastq_files):
    full_text = ""
    for fastq_file in fastq_files:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".fastq", mode='wb') as tmp:
                fastq_file.seek(0)
                tmp.write(fastq_file.read())
                tmp.flush()
                tmp_path = tmp.name
            with open(tmp_path, "r", encoding='utf-8', errors='ignore') as handle:
                sequence_count = 0
                file_text = f"=== FASTQ File: {fastq_file.name} ===\n\n"
                for record in SeqIO.parse(handle, "fastq"):
                    sequence_count += 1
                    # get sequence information
                    sequence_id = record.id
                    description = record.description
                    nucleotide_sequence = str(record.seq)
                    quality_scores = record.letter_annotations.get("phred_quality", [])
                    # basic statistics
                    seq_length = len(nucleotide_sequence)
                    avg_quality = sum(quality_scores) / len(quality_scores) if quality_scores else 0
                    # count nucleotides
                    nucleotide_counts = {
                        'A': nucleotide_sequence.count('A'),
                        'T': nucleotide_sequence.count('T'),
                        'G': nucleotide_sequence.count('G'),
                        'C': nucleotide_sequence.count('C'),
                    }
                    # format sequence information as text for RAG
                    file_text += f"Sequence {sequence_count}:\n"
                    file_text += f"ID: {sequence_id}\n"
                    file_text += f"Description: {description}\n"
                    file_text += f"Length: {seq_length} nucleotides\n"
                    file_text += f"Average Quality Score: {avg_quality:.2f}\n"
                    file_text += f"Nucleotide Composition: A={nucleotide_counts['A']}, T={nucleotide_counts['T']}, G={nucleotide_counts['G']}, C={nucleotide_counts['C']}\n"
                    # for shorter sequences, include the full sequence
                    # for longer sequences, include just the first 100 and last 100 nucleotides
                    if seq_length <= 200:
                        file_text += f"Sequence: {nucleotide_sequence}\n"
                    else:
                        file_text += f"Sequence (first 100): {nucleotide_sequence[:100]}...\n"
                        file_text += f"Sequence (last 100): ...{nucleotide_sequence[-100:]}\n"
                    file_text += f"Quality Scores (first 10): {quality_scores[:10] if quality_scores else 'N/A'}\n"
                    file_text += "\n" + "-" * 50 + "\n\n"
                # add file summary
                file_text += f"File Summary:\n"
                file_text += f"Total sequences in {fastq_file.name}: {sequence_count}\n\n"
                full_text += file_text
            os.remove(tmp_path)
        except Exception as e:
            error_text = f"Error processing FASTQ file {fastq_file.name}: {str(e)}\n\n"
            full_text += error_text
    return full_text.strip()

async def ingest_collection(
    request: Request,
    files: List[UploadFile] = File(...),
    collection_id: str = Form(...),
    collection_name: str = Form(...),
    embedding_model: str = Form("sentence-transformers/all-MiniLM-L6-v2"),
):
    """
    Create a new vectorstore for a specific collection.
    Each collection gets its own isolated vectorstore.
    """
    session_id = request.state.session_id
    
    # Initialize session if needed
    initialize_session(session_id)
    
    # Import USER_DIRS here to avoid circular imports
    from file_handlers.file_tools import USER_DIRS
    
    # Create collection-specific directory
    collection_dir = os.path.join(USER_DIRS, session_id, "collections", collection_id)
    os.makedirs(collection_dir, exist_ok=True)
    
    # Process files and create documents
    all_docs = []
    file_info_list = []
    
    for upload in files:
        ext = Path(upload.filename).suffix.lower()
        metadata = {"source": upload.filename, "filetype": ext}
        
        # Store file info for frontend
        file_info = {
            "name": upload.filename,
            "type": ext.lstrip('.'), 
            "size": upload.size if hasattr(upload, 'size') else 0,
            "dateCreated": time.strftime("%m/%d/%Y")
        }
        file_info_list.append(file_info)
        
        # Handle different file types
        if ext == ".csv":
            df = pd.read_csv(BytesIO(await upload.read()))
            df = clean_dataframe(df)
            for _, row in df.iterrows():
                text = ", ".join(f"{col}: {row[col]}" for col in df.columns)
                doc = Document(page_content=text, metadata=metadata)
                all_docs.append(doc)
                
        elif ext in {".xlsx", ".xls"}:
            df = pd.read_excel(BytesIO(await upload.read()), engine="openpyxl")
            df = clean_dataframe(df)
            for _, row in df.iterrows():
                text = ", ".join(f"{col}: {row[col]}" for col in df.columns)
                doc = Document(page_content=text, metadata=metadata)
                all_docs.append(doc)
                
        elif ext == ".pdf":
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
            tmp.write(await upload.read())
            tmp.flush()
            docs = PyMuPDFLoader(tmp.name).load()
            for doc in docs:
                doc.metadata.update(metadata)
                all_docs.append(doc)
            os.unlink(tmp.name)  # Clean up temp file

        elif ext == ".docx":
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
            tmp.write(await upload.read())
            tmp.flush()
            tmp.close()
            text_content = docx_read([tmp.name])
            if text_content.strip():
                doc = Document(page_content=text_content, metadata=metadata)
                all_docs.append(doc)
            os.unlink(tmp.name)  # Clean up temp file

        elif ext in {".pptx", ".ppt"}:
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            tmp.write(await upload.read())
            tmp.flush()
            tmp.close()
            text_content = pptx_read(tmp.name)
            if text_content.strip():
                doc = Document(page_content=text_content, metadata=metadata)
                all_docs.append(doc)
            os.unlink(tmp.name)
        
        elif ext in {".fastq", ".fq"}:
            text_content = fastq_read([upload])
            if text_content.strip():
                doc = Document(page_content=text_content, metadata=metadata)
                all_docs.append(doc)
            
        elif ext in {".png", ".jpg", ".jpeg", ".gif"}:
            # TODO: implement image embedding with CLIP
            continue
    
    if not all_docs:
        raise HTTPException(400, "No processable files found")
    
    # Create embeddings and vectorstore
    embeddings = HuggingFaceEmbeddings(
        model_name=embedding_model, 
        model_kwargs={'trust_remote_code': True}
    )
    
    # Create text splitter
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=0)
    split_docs = text_splitter.split_documents(all_docs)
    
    # Create vectorstore for this collection
    vectorstore = Chroma.from_documents(
        documents=split_docs,
        embedding=embeddings,
        persist_directory=collection_dir,
    )
    
    # Store collection info in session
    SESSIONS[session_id]["collections"][collection_id] = {
        "vectorstore": vectorstore,
        "name": collection_name,
        "files": file_info_list,
        "created_at": time.time(),
        "embedding_model": embedding_model
    }
    
    # Set as active collection
    SESSIONS[session_id]["active_collection_id"] = collection_id
    
    # Automatically create a chatbot for this collection
    try:
        # Create LLM
        llm = ChatOllama(model="llama3.1", temperature=0)
        
        # Create prompt template
        qa_prompt = PromptTemplate(
            input_variables=["context", "question"],
            template="You are a helpful AI. Use the following context to answer the question:\n\nContext: {context}\n\nQuestion: {question}\n\nAnswer:"
        )
        
        # Create retriever
        retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 2})
        
        # Create chain
        chain = ConversationalRetrievalChain.from_llm(
            llm=llm, 
            retriever=retriever, 
            return_source_documents=True,
            combine_docs_chain_kwargs={"prompt": qa_prompt},
            verbose=True
        )
        
        # Store chain in session
        SESSIONS[session_id]["chain"] = chain
        print(f"DEBUG: Chatbot automatically created for collection {collection_id} in session {session_id}")
        
    except Exception as e:
        print(f"ERROR: Failed to create chatbot automatically: {str(e)}")
        # Don't fail the ingestion if chatbot creation fails
    
    return JSONResponse({
        "session_id": session_id,
        "collection_id": collection_id,
        "collection_name": collection_name,
        "files_processed": len(files),
        "chatbot_created": "chain" in SESSIONS[session_id] and SESSIONS[session_id]["chain"] is not None
    })